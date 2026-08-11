"""Bounded, structured document parsing for the private RAG subsystem."""

from __future__ import annotations

import json
import re
import tomllib
from dataclasses import dataclass
from enum import StrEnum
from pathlib import Path, PurePosixPath
from typing import Any
from zipfile import BadZipFile, ZipFile

from nanobot.rag.config import RagParsingConfig
from nanobot.rag.types import RagErrorCode, SourceKind, SourceLocation


class DocumentFormat(StrEnum):
    PDF = "pdf"
    DOCX = "docx"
    XLSX = "xlsx"
    PPTX = "pptx"
    TEXT = "text"
    MARKDOWN = "markdown"
    CSV = "csv"
    JSON = "json"
    XML = "xml"
    HTML = "html"
    LOG = "log"
    YAML = "yaml"
    TOML = "toml"
    INI = "ini"
    CFG = "cfg"


class ParseCompleteness(StrEnum):
    COMPLETE = "complete"
    TRUNCATED = "truncated"


@dataclass(frozen=True, slots=True)
class ParsedBlock:
    text: str
    location: SourceLocation

    def __post_init__(self) -> None:
        if not self.text.strip():
            raise ValueError("parsed block text must not be empty")


@dataclass(frozen=True, slots=True)
class ParseMetrics:
    file_bytes: int
    total_pages: int | None = None
    processed_pages: int | None = None
    total_slides: int | None = None
    processed_slides: int | None = None
    archive_members: int | None = None
    archive_uncompressed_bytes: int | None = None
    table_rows: int | None = None
    table_cells: int | None = None


@dataclass(frozen=True, slots=True)
class ParsedDocument:
    document_format: DocumentFormat
    blocks: tuple[ParsedBlock, ...]
    completeness: ParseCompleteness
    total_chars: int
    source_chars: int
    metrics: ParseMetrics
    limit_reason: str | None = None

    def __post_init__(self) -> None:
        if not self.blocks:
            raise ValueError("parsed document must contain at least one block")
        if self.completeness is ParseCompleteness.COMPLETE and self.limit_reason is not None:
            raise ValueError("complete parse cannot have a limit reason")
        if self.completeness is ParseCompleteness.TRUNCATED and not self.limit_reason:
            raise ValueError("truncated parse requires a limit reason")


class RagParseError(Exception):
    """A classified parser failure whose message is safe to show to a user."""

    def __init__(self, code: RagErrorCode, safe_message: str) -> None:
        self.code = code
        self.safe_message = safe_message
        super().__init__(safe_message)


_TEXT_FORMATS_BY_SUFFIX = {
    ".txt": DocumentFormat.TEXT,
    ".md": DocumentFormat.MARKDOWN,
    ".markdown": DocumentFormat.MARKDOWN,
    ".csv": DocumentFormat.CSV,
    ".json": DocumentFormat.JSON,
    ".xml": DocumentFormat.XML,
    ".html": DocumentFormat.HTML,
    ".htm": DocumentFormat.HTML,
    ".log": DocumentFormat.LOG,
    ".yaml": DocumentFormat.YAML,
    ".yml": DocumentFormat.YAML,
    ".toml": DocumentFormat.TOML,
    ".ini": DocumentFormat.INI,
    ".cfg": DocumentFormat.CFG,
}
_OFFICE_FORMATS_BY_SUFFIX = {
    ".docx": (DocumentFormat.DOCX, "word/document.xml"),
    ".xlsx": (DocumentFormat.XLSX, "xl/workbook.xml"),
    ".pptx": (DocumentFormat.PPTX, "ppt/presentation.xml"),
}


def detect_document_format(path: str | Path) -> DocumentFormat:
    """Detect a supported format using both suffix and container/content facts."""
    candidate = Path(path)
    try:
        if not candidate.is_file():
            raise RagParseError(RagErrorCode.UNSAFE_DOCUMENT, "无法读取待解析文件")
        suffix = candidate.suffix.lower()
        with candidate.open("rb") as stream:
            header = stream.read(8192)
    except RagParseError:
        raise
    except OSError as exc:
        raise RagParseError(RagErrorCode.UNSAFE_DOCUMENT, "无法读取待解析文件") from exc

    if suffix == ".pdf":
        if not header.startswith(b"%PDF-"):
            raise RagParseError(RagErrorCode.UNSUPPORTED_FORMAT, "文件内容不是有效的 PDF")
        return DocumentFormat.PDF

    if suffix in _OFFICE_FORMATS_BY_SUFFIX:
        expected_format, required_member = _OFFICE_FORMATS_BY_SUFFIX[suffix]
        try:
            with ZipFile(candidate) as archive:
                names = frozenset(archive.namelist())
        except (BadZipFile, OSError) as exc:
            raise RagParseError(
                RagErrorCode.UNSUPPORTED_FORMAT,
                "文件内容不是有效的 Office Open XML 文档",
            ) from exc
        if required_member not in names:
            raise RagParseError(
                RagErrorCode.UNSUPPORTED_FORMAT,
                "文件扩展名与 Office 文档内容不匹配",
            )
        return expected_format

    text_format = _TEXT_FORMATS_BY_SUFFIX.get(suffix)
    if text_format is None:
        raise RagParseError(RagErrorCode.UNSUPPORTED_FORMAT, "不支持此文档格式")
    if not _looks_like_text(header):
        raise RagParseError(
            RagErrorCode.UNSUPPORTED_FORMAT,
            "文件扩展名与二进制内容不匹配",
        )
    return text_format


def parse_document(path: str | Path, config: RagParsingConfig) -> ParsedDocument:
    """Parse one document within configured byte and output bounds."""
    candidate = Path(path)
    document_format = detect_document_format(candidate)
    try:
        file_size = candidate.stat().st_size
    except OSError as exc:
        raise RagParseError(RagErrorCode.UNSAFE_DOCUMENT, "无法检查待解析文件") from exc
    if file_size > config.max_file_bytes:
        raise RagParseError(RagErrorCode.UNSAFE_DOCUMENT, "文件超过允许的大小上限")

    if document_format in _TEXT_FORMATS_BY_SUFFIX.values():
        return _parse_text_document(candidate, document_format, config, file_size)
    if document_format is DocumentFormat.PDF:
        return _parse_pdf_document(candidate, config, file_size)
    if document_format is DocumentFormat.DOCX:
        return _parse_docx_document(candidate, config, file_size)
    if document_format is DocumentFormat.PPTX:
        return _parse_pptx_document(candidate, config, file_size)
    if document_format is DocumentFormat.XLSX:
        return _parse_xlsx_document(candidate, config, file_size)
    raise RagParseError(RagErrorCode.UNSUPPORTED_FORMAT, "此格式的解析器尚不可用")


def _parse_text_document(
    path: Path,
    document_format: DocumentFormat,
    config: RagParsingConfig,
    file_size: int,
) -> ParsedDocument:
    try:
        raw = path.read_bytes()
        text = _decode_text(raw)
    except UnicodeError as exc:
        raise RagParseError(RagErrorCode.UNSUPPORTED_FORMAT, "无法安全识别文本编码") from exc
    except OSError as exc:
        raise RagParseError(RagErrorCode.UNSAFE_DOCUMENT, "无法读取待解析文件") from exc

    _validate_structured_text(text, document_format)

    source_chars = len(text)
    blocks: list[ParsedBlock] = []
    used_chars = 0
    truncated = False
    lines = text.splitlines()
    index = 0
    while index < len(lines):
        while index < len(lines) and not lines[index].strip():
            index += 1
        if index >= len(lines):
            break
        start = index
        paragraph: list[str] = []
        while index < len(lines) and lines[index].strip():
            paragraph.append(lines[index].rstrip())
            index += 1
        block_text = "\n".join(paragraph).strip()
        remaining = config.max_extracted_chars - used_chars
        if remaining <= 0:
            truncated = True
            break
        if len(block_text) > remaining:
            block_text = block_text[:remaining].rstrip()
            truncated = True
        if block_text:
            consumed_lines = block_text.count("\n") + 1
            blocks.append(
                ParsedBlock(
                    text=block_text,
                    location=SourceLocation(
                        kind=SourceKind.TEXT_LINES,
                        line_start=start + 1,
                        line_end=start + consumed_lines,
                    ),
                )
            )
            used_chars += len(block_text)
        if truncated:
            break

    if not blocks:
        raise RagParseError(RagErrorCode.NO_EXTRACTABLE_TEXT, "文档中没有可提取文本")
    completeness = (
        ParseCompleteness.TRUNCATED if truncated else ParseCompleteness.COMPLETE
    )
    return ParsedDocument(
        document_format=document_format,
        blocks=tuple(blocks),
        completeness=completeness,
        total_chars=used_chars,
        source_chars=source_chars,
        metrics=ParseMetrics(file_bytes=file_size),
        limit_reason="max_extracted_chars" if truncated else None,
    )


def _parse_pdf_document(
    path: Path,
    config: RagParsingConfig,
    file_size: int,
) -> ParsedDocument:
    from pypdf import PdfReader
    from pypdf.errors import FileNotDecryptedError, PdfReadError

    try:
        reader = PdfReader(path, strict=False)
        if reader.is_encrypted:
            raise RagParseError(
                RagErrorCode.ENCRYPTED_DOCUMENT,
                "不支持加密 PDF 文档",
            )
        total_pages = len(reader.pages)
        page_limit = min(total_pages, config.max_pdf_pages)
        blocks: list[ParsedBlock] = []
        used_chars = 0
        source_chars = 0
        processed_pages = 0
        limit_reason = "max_pdf_pages" if page_limit < total_pages else None

        for page_index in range(page_limit):
            page = reader.pages[page_index]
            contents = page.get_contents()
            if contents is not None:
                stream_bytes = len(contents.get_data())
                if stream_bytes > config.max_pdf_content_stream_bytes:
                    raise RagParseError(
                        RagErrorCode.UNSAFE_DOCUMENT,
                        "PDF 页面内容流超过安全上限",
                    )
            page_text = (page.extract_text() or "").strip()
            processed_pages += 1
            source_chars += len(page_text)
            if not page_text:
                continue
            remaining = config.max_extracted_chars - used_chars
            if remaining <= 0:
                limit_reason = "max_extracted_chars"
                break
            if len(page_text) > remaining:
                page_text = page_text[:remaining].rstrip()
                limit_reason = "max_extracted_chars"
            if page_text:
                blocks.append(
                    ParsedBlock(
                        text=page_text,
                        location=SourceLocation(
                            kind=SourceKind.PDF_PAGE,
                            page=page_index + 1,
                        ),
                    )
                )
                used_chars += len(page_text)
            if limit_reason == "max_extracted_chars":
                break
    except RagParseError:
        raise
    except (FileNotDecryptedError, PdfReadError) as exc:
        raise RagParseError(RagErrorCode.UNSAFE_DOCUMENT, "PDF 文档结构无效") from exc
    except OSError as exc:
        raise RagParseError(RagErrorCode.UNSAFE_DOCUMENT, "无法读取 PDF 文档") from exc
    except Exception as exc:
        raise RagParseError(RagErrorCode.UNSAFE_DOCUMENT, "PDF 文档解析失败") from exc

    if not blocks:
        raise RagParseError(RagErrorCode.NO_EXTRACTABLE_TEXT, "PDF 中没有可提取文本；首版不执行 OCR")
    completeness = (
        ParseCompleteness.TRUNCATED
        if limit_reason is not None
        else ParseCompleteness.COMPLETE
    )
    return ParsedDocument(
        document_format=DocumentFormat.PDF,
        blocks=tuple(blocks),
        completeness=completeness,
        total_chars=used_chars,
        source_chars=source_chars,
        metrics=ParseMetrics(
            file_bytes=file_size,
            total_pages=total_pages,
            processed_pages=processed_pages,
        ),
        limit_reason=limit_reason,
    )


@dataclass(frozen=True, slots=True)
class _ArchiveMetrics:
    members: int
    uncompressed_bytes: int


def _inspect_office_archive(path: Path, config: RagParsingConfig) -> _ArchiveMetrics:
    try:
        with ZipFile(path) as archive:
            members = archive.infolist()
    except (BadZipFile, OSError) as exc:
        raise RagParseError(RagErrorCode.UNSAFE_DOCUMENT, "Office 文档容器无效") from exc
    if len(members) > config.max_archive_members:
        raise RagParseError(RagErrorCode.UNSAFE_DOCUMENT, "Office 文档包含过多内部文件")
    total_size = 0
    for member in members:
        member_path = PurePosixPath(member.filename)
        if member.flag_bits & 0x1:
            raise RagParseError(
                RagErrorCode.ENCRYPTED_DOCUMENT,
                "不支持加密 Office 文档",
            )
        if member_path.is_absolute() or ".." in member_path.parts:
            raise RagParseError(RagErrorCode.UNSAFE_DOCUMENT, "Office 文档包含不安全路径")
        if member.file_size > config.max_archive_member_bytes:
            raise RagParseError(RagErrorCode.UNSAFE_DOCUMENT, "Office 内部文件超过安全上限")
        if member.filename.lower().endswith("vbaproject.bin"):
            raise RagParseError(RagErrorCode.UNSAFE_DOCUMENT, "不支持包含宏的 Office 文档")
        total_size += member.file_size
        if total_size > config.max_archive_uncompressed_bytes:
            raise RagParseError(RagErrorCode.UNSAFE_DOCUMENT, "Office 文档解压后超过安全上限")
    return _ArchiveMetrics(len(members), total_size)


def _append_bounded_block(
    blocks: list[ParsedBlock],
    *,
    text: str,
    location: SourceLocation,
    used_chars: int,
    limit: int,
) -> tuple[int, bool]:
    normalized = text.strip()
    if not normalized:
        return used_chars, False
    remaining = limit - used_chars
    if remaining <= 0:
        return used_chars, True
    truncated = len(normalized) > remaining
    if truncated:
        normalized = normalized[:remaining].rstrip()
    if normalized:
        blocks.append(ParsedBlock(text=normalized, location=location))
        used_chars += len(normalized)
    return used_chars, truncated


def _heading_level(style_name: str) -> int | None:
    match = re.fullmatch(r"Heading\s+([1-9])", style_name, flags=re.IGNORECASE)
    return int(match.group(1)) if match else None


def _parse_docx_document(
    path: Path,
    config: RagParsingConfig,
    file_size: int,
) -> ParsedDocument:
    from docx import Document
    from docx.text.paragraph import Paragraph

    archive = _inspect_office_archive(path, config)
    blocks: list[ParsedBlock] = []
    heading_path: list[str] = []
    used_chars = 0
    source_chars = 0
    table_rows = 0
    table_cells = 0
    truncated = False
    try:
        document = Document(str(path))
        for block in document.iter_inner_content():
            if isinstance(block, Paragraph):
                text = block.text.strip()
                if not text:
                    continue
                source_chars += len(text)
                level = _heading_level(getattr(block.style, "name", ""))
                if level is not None:
                    heading_path[level - 1 :] = [text]
                location = (
                    SourceLocation(kind=SourceKind.HEADING, heading_path=tuple(heading_path))
                    if heading_path
                    else SourceLocation(kind=SourceKind.UNKNOWN)
                )
                used_chars, truncated = _append_bounded_block(
                    blocks,
                    text=text,
                    location=location,
                    used_chars=used_chars,
                    limit=config.max_extracted_chars,
                )
                if truncated:
                    break
                continue
            for row in block.rows:
                table_rows += 1
                if table_rows > config.max_table_rows:
                    raise RagParseError(RagErrorCode.UNSAFE_DOCUMENT, "DOCX 表格行数超过安全上限")
                cells = [cell.text.strip() for cell in row.cells]
                table_cells += len(cells)
                if table_cells > config.max_table_cells:
                    raise RagParseError(RagErrorCode.UNSAFE_DOCUMENT, "DOCX 表格单元格超过安全上限")
                row_text = "\t".join(cells).strip()
                source_chars += len(row_text)
                location = (
                    SourceLocation(kind=SourceKind.HEADING, heading_path=tuple(heading_path))
                    if heading_path
                    else SourceLocation(kind=SourceKind.UNKNOWN)
                )
                used_chars, truncated = _append_bounded_block(
                    blocks,
                    text=row_text,
                    location=location,
                    used_chars=used_chars,
                    limit=config.max_extracted_chars,
                )
                if truncated:
                    break
            if truncated:
                break
    except RagParseError:
        raise
    except Exception as exc:
        raise RagParseError(RagErrorCode.UNSAFE_DOCUMENT, "DOCX 文档解析失败") from exc

    if not blocks:
        raise RagParseError(RagErrorCode.NO_EXTRACTABLE_TEXT, "DOCX 中没有可提取文本；首版不执行 OCR")
    return ParsedDocument(
        document_format=DocumentFormat.DOCX,
        blocks=tuple(blocks),
        completeness=(
            ParseCompleteness.TRUNCATED if truncated else ParseCompleteness.COMPLETE
        ),
        total_chars=used_chars,
        source_chars=source_chars,
        metrics=ParseMetrics(
            file_bytes=file_size,
            archive_members=archive.members,
            archive_uncompressed_bytes=archive.uncompressed_bytes,
            table_rows=table_rows,
            table_cells=table_cells,
        ),
        limit_reason="max_extracted_chars" if truncated else None,
    )


def _collect_pptx_text(
    shape: Any,
    *,
    depth: int,
    config: RagParsingConfig,
    counters: list[int],
) -> list[str]:
    if depth > config.max_structure_depth:
        raise RagParseError(RagErrorCode.UNSAFE_DOCUMENT, "PPTX 组合对象嵌套超过安全上限")
    sub_shapes = getattr(shape, "shapes", None)
    if sub_shapes is not None:
        result: list[str] = []
        for sub_shape in sub_shapes:
            result.extend(
                _collect_pptx_text(
                    sub_shape,
                    depth=depth + 1,
                    config=config,
                    counters=counters,
                )
            )
        return result
    if getattr(shape, "has_table", False):
        result = []
        for row in shape.table.rows:
            counters[0] += 1
            if counters[0] > config.max_table_rows:
                raise RagParseError(RagErrorCode.UNSAFE_DOCUMENT, "PPTX 表格行数超过安全上限")
            cells = [cell.text.strip() for cell in row.cells]
            counters[1] += len(cells)
            if counters[1] > config.max_table_cells:
                raise RagParseError(RagErrorCode.UNSAFE_DOCUMENT, "PPTX 表格单元格超过安全上限")
            if row_text := "\t".join(cells).strip():
                result.append(row_text)
        return result
    text = str(getattr(shape, "text", "")).strip()
    return [text] if text else []


def _parse_pptx_document(
    path: Path,
    config: RagParsingConfig,
    file_size: int,
) -> ParsedDocument:
    from pptx import Presentation

    archive = _inspect_office_archive(path, config)
    blocks: list[ParsedBlock] = []
    used_chars = 0
    source_chars = 0
    counters = [0, 0]
    limit_reason: str | None = None
    try:
        presentation = Presentation(str(path))
        total_slides = len(presentation.slides)
        slide_limit = min(total_slides, config.max_presentation_slides)
        if slide_limit < total_slides:
            limit_reason = "max_presentation_slides"
        processed_slides = 0
        for slide_index in range(slide_limit):
            slide = presentation.slides[slide_index]
            processed_slides += 1
            for shape in slide.shapes:
                for text in _collect_pptx_text(
                    shape,
                    depth=1,
                    config=config,
                    counters=counters,
                ):
                    source_chars += len(text)
                    used_chars, truncated = _append_bounded_block(
                        blocks,
                        text=text,
                        location=SourceLocation(
                            kind=SourceKind.SLIDE,
                            slide=slide_index + 1,
                        ),
                        used_chars=used_chars,
                        limit=config.max_extracted_chars,
                    )
                    if truncated:
                        limit_reason = "max_extracted_chars"
                        break
                if limit_reason == "max_extracted_chars":
                    break
            if limit_reason == "max_extracted_chars":
                break
    except RagParseError:
        raise
    except Exception as exc:
        raise RagParseError(RagErrorCode.UNSAFE_DOCUMENT, "PPTX 文档解析失败") from exc

    if not blocks:
        raise RagParseError(RagErrorCode.NO_EXTRACTABLE_TEXT, "PPTX 中没有可提取文本；首版不执行 OCR")
    return ParsedDocument(
        document_format=DocumentFormat.PPTX,
        blocks=tuple(blocks),
        completeness=(
            ParseCompleteness.TRUNCATED
            if limit_reason is not None
            else ParseCompleteness.COMPLETE
        ),
        total_chars=used_chars,
        source_chars=source_chars,
        metrics=ParseMetrics(
            file_bytes=file_size,
            total_slides=total_slides,
            processed_slides=processed_slides,
            archive_members=archive.members,
            archive_uncompressed_bytes=archive.uncompressed_bytes,
            table_rows=counters[0],
            table_cells=counters[1],
        ),
        limit_reason=limit_reason,
    )


def _parse_xlsx_document(
    path: Path,
    config: RagParsingConfig,
    file_size: int,
) -> ParsedDocument:
    from openpyxl import load_workbook

    archive = _inspect_office_archive(path, config)
    blocks: list[ParsedBlock] = []
    used_chars = 0
    source_chars = 0
    table_rows = 0
    table_cells = 0
    truncated = False
    try:
        workbook = load_workbook(
            path,
            read_only=True,
            data_only=True,
            keep_links=False,
        )
        try:
            if len(workbook.sheetnames) > config.max_spreadsheet_sheets:
                raise RagParseError(
                    RagErrorCode.UNSAFE_DOCUMENT,
                    "XLSX 工作表数量超过安全上限",
                )
            for sheet_name in workbook.sheetnames:
                worksheet = workbook[sheet_name]
                for row_number, row in enumerate(
                    worksheet.iter_rows(values_only=True),
                    start=1,
                ):
                    table_rows += 1
                    if table_rows > config.max_table_rows:
                        raise RagParseError(
                            RagErrorCode.UNSAFE_DOCUMENT,
                            "XLSX 行数超过安全上限",
                        )
                    table_cells += len(row)
                    if table_cells > config.max_table_cells:
                        raise RagParseError(
                            RagErrorCode.UNSAFE_DOCUMENT,
                            "XLSX 单元格数量超过安全上限",
                        )
                    row_text = "\t".join(
                        str(value) if value is not None else "" for value in row
                    ).rstrip()
                    if not row_text.strip():
                        continue
                    source_chars += len(row_text)
                    used_chars, truncated = _append_bounded_block(
                        blocks,
                        text=row_text,
                        location=SourceLocation(
                            kind=SourceKind.SPREADSHEET_ROWS,
                            sheet=sheet_name,
                            row_start=row_number,
                            row_end=row_number,
                        ),
                        used_chars=used_chars,
                        limit=config.max_extracted_chars,
                    )
                    if truncated:
                        break
                if truncated:
                    break
        finally:
            workbook.close()
    except RagParseError:
        raise
    except Exception as exc:
        raise RagParseError(RagErrorCode.UNSAFE_DOCUMENT, "XLSX 文档解析失败") from exc

    if not blocks:
        raise RagParseError(RagErrorCode.NO_EXTRACTABLE_TEXT, "XLSX 中没有可提取文本")
    return ParsedDocument(
        document_format=DocumentFormat.XLSX,
        blocks=tuple(blocks),
        completeness=(
            ParseCompleteness.TRUNCATED if truncated else ParseCompleteness.COMPLETE
        ),
        total_chars=used_chars,
        source_chars=source_chars,
        metrics=ParseMetrics(
            file_bytes=file_size,
            archive_members=archive.members,
            archive_uncompressed_bytes=archive.uncompressed_bytes,
            table_rows=table_rows,
            table_cells=table_cells,
        ),
        limit_reason="max_extracted_chars" if truncated else None,
    )


def _looks_like_text(header: bytes) -> bool:
    if not header:
        return True
    if header.startswith(
        (
            b"%PDF-",
            b"PK\x03\x04",
            b"\x89PNG\r\n\x1a\n",
            b"GIF87a",
            b"GIF89a",
            b"\xff\xd8\xff",
            b"\xd0\xcf\x11\xe0\xa1\xb1\x1a\xe1",
        )
    ):
        return False
    if header.startswith((b"\xff\xfe", b"\xfe\xff")):
        return True
    if b"\x00" in header:
        return False
    control_count = sum(byte < 9 or 13 < byte < 32 for byte in header)
    return control_count / len(header) < 0.02


def _decode_text(raw: bytes) -> str:
    if raw.startswith((b"\xff\xfe", b"\xfe\xff")):
        return raw.decode("utf-16")
    try:
        return raw.decode("utf-8-sig")
    except UnicodeDecodeError:
        # Latin-1 is deterministic and lossless; the content heuristic already
        # rejected binary control bytes before this fallback is reached.
        return raw.decode("latin-1")


def _validate_structured_text(text: str, document_format: DocumentFormat) -> None:
    if document_format is DocumentFormat.JSON:
        try:
            json.loads(text)
        except (json.JSONDecodeError, RecursionError) as exc:
            raise RagParseError(RagErrorCode.UNSUPPORTED_FORMAT, "JSON 内容无效") from exc
        return
    if document_format is DocumentFormat.XML:
        if re.search(r"<!\s*(?:DOCTYPE|ENTITY)\b", text, flags=re.IGNORECASE):
            raise RagParseError(RagErrorCode.UNSAFE_DOCUMENT, "XML 包含不允许的实体或文档类型声明")
        try:
            from defusedxml.ElementTree import fromstring

            fromstring(text)
        except Exception as exc:
            raise RagParseError(RagErrorCode.UNSUPPORTED_FORMAT, "XML 内容无效") from exc
        return
    if document_format is DocumentFormat.HTML:
        if not re.search(
            r"<\s*(?:!doctype\s+html|html|head|body|p|div|h[1-6]|table|ul|ol|article|section)\b",
            text,
            flags=re.IGNORECASE,
        ):
            raise RagParseError(RagErrorCode.UNSUPPORTED_FORMAT, "HTML 内容无效")
        return
    if document_format is DocumentFormat.TOML:
        try:
            tomllib.loads(text)
        except tomllib.TOMLDecodeError as exc:
            raise RagParseError(RagErrorCode.UNSUPPORTED_FORMAT, "TOML 内容无效") from exc


__all__ = [
    "DocumentFormat",
    "ParseCompleteness",
    "ParseMetrics",
    "ParsedBlock",
    "ParsedDocument",
    "RagParseError",
    "detect_document_format",
    "parse_document",
]

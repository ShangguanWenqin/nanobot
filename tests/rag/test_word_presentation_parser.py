from __future__ import annotations

from pathlib import Path
from zipfile import ZIP_DEFLATED, ZipFile

import pytest
from docx import Document
from pptx import Presentation
from pptx.util import Inches

from nanobot.rag.config import RagParsingConfig
from nanobot.rag.parser import ParseCompleteness, RagParseError, parse_document
from nanobot.rag.types import RagErrorCode, SourceKind


def _mark_zip_encrypted(path: Path) -> None:
    data = bytearray(path.read_bytes())
    cursor = 0
    while cursor < len(data) - 10:
        if data[cursor : cursor + 4] == b"PK\x03\x04":
            flags = int.from_bytes(data[cursor + 6 : cursor + 8], "little") | 1
            data[cursor + 6 : cursor + 8] = flags.to_bytes(2, "little")
        elif data[cursor : cursor + 4] == b"PK\x01\x02":
            flags = int.from_bytes(data[cursor + 8 : cursor + 10], "little") | 1
            data[cursor + 8 : cursor + 10] = flags.to_bytes(2, "little")
        cursor += 1
    path.write_bytes(data)


def test_docx_preserves_heading_path_and_table_rows(tmp_path: Path) -> None:
    path = tmp_path / "guide.docx"
    document = Document()
    document.add_heading("部署", level=1)
    document.add_paragraph("准备运行环境。")
    table = document.add_table(rows=2, cols=2)
    table.cell(0, 0).text = "名称"
    table.cell(0, 1).text = "值"
    table.cell(1, 0).text = "模式"
    table.cell(1, 1).text = "本地"
    document.add_heading("验证", level=2)
    document.add_paragraph("运行测试。")
    document.save(path)

    result = parse_document(path, RagParsingConfig())

    assert [block.text for block in result.blocks] == [
        "部署",
        "准备运行环境。",
        "名称\t值",
        "模式\t本地",
        "验证",
        "运行测试。",
    ]
    assert result.blocks[1].location.kind is SourceKind.HEADING
    assert result.blocks[1].location.heading_path == ("部署",)
    assert result.blocks[-1].location.heading_path == ("部署", "验证")
    assert result.metrics.table_rows == 2
    assert result.metrics.table_cells == 4


def test_docx_table_limit_is_a_safe_error(tmp_path: Path) -> None:
    path = tmp_path / "table.docx"
    document = Document()
    document.add_table(rows=2, cols=2)
    document.save(path)

    with pytest.raises(RagParseError) as exc_info:
        parse_document(path, RagParsingConfig(max_table_rows=1))

    assert exc_info.value.code is RagErrorCode.UNSAFE_DOCUMENT


def test_pptx_recurses_into_group_shapes_and_preserves_slide(tmp_path: Path) -> None:
    path = tmp_path / "deck.pptx"
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    group = slide.shapes.add_group_shape()
    box = group.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(1))
    box.text = "分组内文本"
    table_shape = slide.shapes.add_table(1, 2, Inches(1), Inches(2), Inches(4), Inches(1))
    table_shape.table.cell(0, 0).text = "A"
    table_shape.table.cell(0, 1).text = "B"
    presentation.save(path)

    result = parse_document(path, RagParsingConfig())

    assert [block.text for block in result.blocks] == ["分组内文本", "A\tB"]
    assert all(block.location.kind is SourceKind.SLIDE for block in result.blocks)
    assert all(block.location.slide == 1 for block in result.blocks)


def test_office_macro_payload_is_rejected(tmp_path: Path) -> None:
    path = tmp_path / "macro.docx"
    document = Document()
    document.add_paragraph("normal text")
    document.save(path)
    with ZipFile(path, "a", compression=ZIP_DEFLATED) as archive:
        archive.writestr("word/vbaProject.bin", b"macro")

    with pytest.raises(RagParseError) as exc_info:
        parse_document(path, RagParsingConfig())

    assert exc_info.value.code is RagErrorCode.UNSAFE_DOCUMENT


def test_encrypted_office_archive_is_rejected(tmp_path: Path) -> None:
    path = tmp_path / "encrypted.docx"
    document = Document()
    document.add_paragraph("secret")
    document.save(path)
    _mark_zip_encrypted(path)

    with pytest.raises(RagParseError) as exc_info:
        parse_document(path, RagParsingConfig())

    assert exc_info.value.code is RagErrorCode.ENCRYPTED_DOCUMENT


def test_presentation_slide_limit_is_explicit(tmp_path: Path) -> None:
    path = tmp_path / "deck.pptx"
    presentation = Presentation()
    for text in ("One", "Two"):
        slide = presentation.slides.add_slide(presentation.slide_layouts[6])
        slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(1)).text = text
    presentation.save(path)

    result = parse_document(path, RagParsingConfig(max_presentation_slides=1))

    assert [block.text for block in result.blocks] == ["One"]
    assert result.completeness is ParseCompleteness.TRUNCATED
    assert result.limit_reason == "max_presentation_slides"


def test_image_only_presentation_is_not_ocrd(tmp_path: Path) -> None:
    path = tmp_path / "empty.pptx"
    presentation = Presentation()
    presentation.slides.add_slide(presentation.slide_layouts[6])
    presentation.save(path)

    with pytest.raises(RagParseError) as exc_info:
        parse_document(path, RagParsingConfig())

    assert exc_info.value.code is RagErrorCode.NO_EXTRACTABLE_TEXT
